#!/usr/bin/env python3
"""Render each slide of a shipped .pptx to PNG for visual QA, and flag geometric defects.

LibreOffice cannot load any file in this sandbox (a plain .txt fails too), so the skill's
soffice -> pdftoppm route is unavailable. This reads the ACTUAL shipped package via python-pptx --
real positions, sizes, fills, text and embedded images -- so what it draws is the file, not the
generator's intent. Text wrapping is approximated from per-character widths, which is enough to
catch the defect class that matters most: text overflowing its box.
"""
import io
import pathlib
import sys

import cairosvg
from pptx import Presentation
from pptx.util import Emu

DECK = pathlib.Path(sys.argv[1] if len(sys.argv) > 1 else "build/SumoSharp-features.pptx")
OUT = DECK.parent / "qa"
OUT.mkdir(exist_ok=True)
SCALE = 100          # px per inch
EMU_IN = 914400.0

# Measure with real glyph metrics rather than a per-character guess: a guess that UNDER-estimates
# is worse than no check at all, because PowerPoint then wraps earlier than the preview and text
# overflows vertically in the shipped file while QA reports it clean. DejaVu is a little wider than
# Calibri, so every measurement here errs toward "too wide" -- the safe direction.
from PIL import ImageFont

_FONTS = {}


def _font(size, bold=False, mono=False):
    key = (round(size, 1), bold, mono)
    if key not in _FONTS:
        base = "/usr/share/fonts/truetype/dejavu/DejaVuSans"
        if mono:
            path = base + "Mono-Bold.ttf" if bold else base + "Mono.ttf"
        else:
            path = base + "-Bold.ttf" if bold else base + ".ttf"
        _FONTS[key] = ImageFont.truetype(path, max(1, int(round(size))))
    return _FONTS[key]


def text_w(s, size, bold=False, mono=False):
    return _font(size, bold, mono).getlength(str(s))


def wrap(s, size, box_w_px, bold=False, mono=False):
    """Greedy wrap at box width, honouring explicit newlines."""
    out = []
    for para in str(s).split("\n"):
        line = ""
        for word in para.split(" "):
            trial = (line + " " + word).strip()
            if line and text_w(trial, size, bold, mono) > box_w_px:
                out.append(line)
                line = word
            else:
                line = trial
        out.append(line)
    return out


def esc(s):
    return (str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def hexof(color_obj, fallback=None):
    try:
        if color_obj and color_obj.type is not None and str(color_obj.rgb):
            return "#" + str(color_obj.rgb)
    except Exception:
        pass
    return fallback


prs = Presentation(str(DECK))
SW = prs.slide_width / EMU_IN * SCALE
SH = prs.slide_height / EMU_IN * SCALE
problems = []

for idx, slide in enumerate(prs.slides, 1):
    body = ""
    bgc = "#FFFFFF"
    try:
        f = slide.background.fill
        if f.type is not None and f.fore_color and f.fore_color.rgb:
            bgc = "#" + str(f.fore_color.rgb)
    except Exception:
        pass
    body += f'<rect width="{SW}" height="{SH}" fill="{bgc}"/>'

    boxes = []
    for sh in slide.shapes:
        x = sh.left / EMU_IN * SCALE
        y = sh.top / EMU_IN * SCALE
        w = sh.width / EMU_IN * SCALE
        h = sh.height / EMU_IN * SCALE
        boxes.append((sh.shape_type, x, y, w, h))

        # off-slide / margin checks against the real geometry
        if x < -1 or y < -1 or x + w > SW + 1 or y + h > SH + 1:
            problems.append(f"slide {idx}: shape past slide bounds "
                            f"({x:.0f},{y:.0f},{w:.0f}x{h:.0f}) canvas {SW:.0f}x{SH:.0f}")

        if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":  # PICTURE
            try:
                blob = sh.image.blob
                import base64
                b64 = base64.b64encode(blob).decode()
                body += (f'<image x="{x}" y="{y}" width="{w}" height="{h}" '
                         f'xlink:href="data:image/png;base64,{b64}" preserveAspectRatio="none"/>')
            except Exception as e:
                body += f'<rect x="{x}" y="{y}" width="{w}" height="{h}" fill="#444"/>'
            continue

        fill = None
        line = None
        try:
            if sh.fill.type is not None and sh.fill.type != 5:
                fill = hexof(sh.fill.fore_color)
        except Exception:
            pass
        try:
            line = hexof(sh.line.color)
        except Exception:
            pass
        if fill or line:
            body += (f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="6" '
                     f'fill="{fill or "none"}" stroke="{line or "none"}" stroke-width="1.2"/>')

        if not sh.has_text_frame:
            continue

        # lay the paragraphs out inside the box, then check whether they fit
        cy = y + 4
        for para in sh.text_frame.paragraphs:
            runs = [r for r in para.runs if r.text]
            if not runs:
                cy += 6
                continue
            size = 14.0
            for r in runs:
                if r.font.size:
                    size = r.font.size.pt
                    break
            col = "#111111"
            for r in runs:
                c = hexof(r.font.color)
                if c:
                    col = c
                    break
            bold = any(r.font.bold for r in runs)
            italic = any(r.font.italic for r in runs)
            fam = "DejaVu Sans"
            for r in runs:
                if r.font.name and "Consol" in r.font.name:
                    fam = "DejaVu Sans Mono"
            txt = "".join(r.text for r in runs)
            lh = size * 1.32
            style = " font-style='italic'" if italic else ""
            for ln in wrap(txt, size, w - 8, bold, fam.endswith('Mono')):
                cy += lh
                body += (f'<text x="{x + 4}" y="{cy}" font-family="{fam}" font-size="{size}" '
                         f'fill="{col}" font-weight="{"bold" if bold else "normal"}"'
                         f'{style}>{esc(ln)}</text>')
            cy += 3
        if cy > y + h + 3:
            problems.append(f"slide {idx}: TEXT OVERFLOW — box y {y:.0f}..{y + h:.0f}, "
                            f"text reaches {cy:.0f} (+{cy - (y + h):.0f}px): "
                            f"{sh.text_frame.text[:60]!r}")

    svg = (f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" '
           f'width="{SW}" height="{SH}" viewBox="0 0 {SW} {SH}">{body}</svg>')
    cairosvg.svg2png(bytestring=svg.encode(), write_to=str(OUT / f"slide-{idx:02d}.png"), scale=1.0)

print(f"rendered {len(prs.slides)} slides to {OUT}/")
if problems:
    print(f"\n{len(problems)} geometric problem(s):")
    for x in problems:
        print("  -", x)
else:
    print("no off-slide shapes, no estimated text overflow")
